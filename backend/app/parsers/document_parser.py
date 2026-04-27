import io
import json
from pathlib import Path

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".md",
    ".txt",
    ".pdf",
    ".docx",
    ".csv",
    ".json",
    ".html",
    ".xml",
    ".xlsx",
    ".pptx",
}


def parse_document(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {suffix}")

    if suffix in {".md", ".txt"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".csv":
        return _parse_csv(path)
    if suffix == ".json":
        return _parse_json(path)
    if suffix in {".html", ".xml"}:
        return _parse_markup(path)
    if suffix == ".xlsx":
        return _parse_xlsx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)

    raise ValueError(f"Unsupported file type: {suffix}")


def _parse_pdf(path: Path) -> str:
    import fitz
    import pymupdf4llm

    document = fitz.open(path)
    pages: list[str] = []

    try:
        for page_number, page in enumerate(document):
            page_markdown = pymupdf4llm.to_markdown(
                path.as_posix(),
                pages=[page_number],
            ).strip()

            if _has_low_text(page_markdown):
                ocr_text = _extract_ocr_text(page)
                if ocr_text:
                    page_markdown = f"{page_markdown}\n\n{ocr_text}".strip()

            if page_markdown:
                pages.append(page_markdown)
    finally:
        document.close()

    return "\n\n".join(pages)


def _has_low_text(text: str, minimum_characters: int = 40) -> bool:
    condensed = " ".join(text.split())
    return len(condensed) < minimum_characters


def _extract_ocr_text(page) -> str:
    try:
        import cv2
        import numpy as np
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        print("OCR dependencies are not installed. Skipping OCR extraction.")
        return ""

    try:
        pixmap = page.get_pixmap(matrix=page.derotation_matrix)
        image_bytes = pixmap.tobytes("png")
        image_array = np.frombuffer(image_bytes, dtype=np.uint8)
        image = cv2.imdecode(image_array, cv2.IMREAD_COLOR)
        if image is None:
            return ""

        engine = RapidOCR()
        result, _ = engine(image)
        if not result:
            return ""

        lines = [item[1] for item in result if len(item) > 1 and item[1]]
        return "\n".join(lines).strip()
    except Exception:
        print("An error occurred while performing OCR.")
        return ""


def _parse_docx(path: Path) -> str:
    document = DocxDocument(path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_csv(path: Path) -> str:
    dataframe = pd.read_csv(path)
    return dataframe.to_markdown(index=False)


def _parse_json(path: Path) -> str:
    payload = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(payload, indent=2, ensure_ascii=False)


def _parse_markup(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(content, "html.parser")
    return soup.get_text(separator="\n", strip=True)


def _parse_xlsx(path: Path) -> str:
    workbook = pd.read_excel(path, sheet_name=None)
    sections: list[str] = []
    for sheet_name, dataframe in workbook.items():
        sections.append(f"## Sheet: {sheet_name}")
        sections.append(dataframe.fillna("").to_markdown(index=False))
    return "\n\n".join(sections)


def _parse_pptx(path: Path) -> str:
    presentation = Presentation(path)
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"## Slide {index}\n\n" + "\n".join(texts))

    return "\n\n".join(slides)
